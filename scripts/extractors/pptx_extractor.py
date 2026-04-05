from __future__ import annotations

from dataclasses import asdict, dataclass
from io import BytesIO
from pathlib import Path
from typing import Any

from PIL import Image

from core.chunking import clean_text, make_sort_key
from image_utils.caption import find_best_caption_for_image
from image_utils.filtering import is_diagram_image
from image_utils.ocr import compute_ocr

try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_SHAPE_TYPE
except Exception:  # pragma: no cover
    Presentation = None
    MSO_AUTO_SHAPE_TYPE = None
    MSO_SHAPE_TYPE = None


@dataclass
class TextBlock:
    block_id: str
    page_number: int
    block_index: int
    bbox: dict[str, float]
    text: str
    sort_key: str
    document_order: int
    shape_name: str | None = None
    shape_type: str | None = None


@dataclass
class ImageItem:
    image_index: int
    page_number: int
    block_index: int
    bbox: dict[str, float]
    image_path: str
    ext: str
    caption_text: str
    filter_reason: str
    ocr_text: str
    ocr_word_count: int
    ocr_avg_conf: float
    sort_key: str
    document_order: int
    shape_name: str | None = None


@dataclass
class ExtractedPPTX:
    source_path: str
    file_type: str
    slide_count: int
    text_blocks: list[TextBlock]
    images: list[ImageItem]
    stats: dict[str, Any]


@dataclass
class _Rect:
    x0: float
    y0: float
    x1: float
    y1: float


def _save_image(image_bytes: bytes, image_path: Path) -> None:
    image_path.parent.mkdir(parents=True, exist_ok=True)
    image_path.write_bytes(image_bytes)


def _shape_bbox(shape: Any) -> dict[str, float]:
    left = float(getattr(shape, "left", 0.0))
    top = float(getattr(shape, "top", 0.0))
    width = float(getattr(shape, "width", 0.0))
    height = float(getattr(shape, "height", 0.0))
    return {
        "x0": left,
        "y0": top,
        "x1": left + width,
        "y1": top + height,
    }


def _shape_rect(shape: Any) -> _Rect:
    bbox = _shape_bbox(shape)
    return _Rect(
        x0=float(bbox["x0"]),
        y0=float(bbox["y0"]),
        x1=float(bbox["x1"]),
        y1=float(bbox["y1"]),
    )


def _shape_type_name(shape: Any) -> str:
    shape_type = getattr(shape, "shape_type", None)
    if shape_type is not None and hasattr(shape_type, "name"):
        return str(shape_type.name)
    return str(shape_type or "unknown")


def _auto_shape_name(shape: Any) -> str | None:
    try:
        auto_shape = getattr(shape, "auto_shape_type", None)
    except Exception:
        return None
    if auto_shape is not None and hasattr(auto_shape, "name"):
        return str(auto_shape.name)
    return str(auto_shape) if auto_shape is not None else None


def _shape_text(shape: Any) -> str:
    try:
        return clean_text(str(getattr(shape, "text", "") or ""))
    except Exception:
        return ""


def _table_to_text(shape: Any, slide_number: int, table_index: int) -> str:
    rows_out: list[str] = [f"Slide {slide_number} Table {table_index}"]
    table = getattr(shape, "table", None)
    if table is None:
        return ""
    for ridx, row in enumerate(table.rows, 1):
        cells: list[str] = []
        for cidx, cell in enumerate(row.cells, 1):
            cell_text = clean_text(cell.text)
            if cell_text:
                cells.append(f"col_{cidx}: {cell_text}")
        if cells:
            rows_out.append(f"row_{ridx}: " + " ; ".join(cells))
    return clean_text("\n".join(rows_out))


def _flatten_shapes(shapes: Any) -> list[Any]:
    out: list[Any] = []
    for shape in shapes:
        if getattr(shape, "shape_type", None) == getattr(MSO_SHAPE_TYPE, "GROUP", None) and hasattr(shape, "shapes"):
            out.extend(_flatten_shapes(shape.shapes))
        else:
            out.append(shape)
    return out


def _auto_shape_category(auto_shape_name: str | None) -> str:
    """Classify an auto_shape_type name into a semantic role for diagram descriptions."""
    if not auto_shape_name:
        return "SHAPE"
    n = auto_shape_name.upper()
    if "DIAMOND" in n:
        return "DECISION"
    if any(k in n for k in ("OVAL", "ELLIPSE", "TERMINATOR")):
        return "TERMINAL"
    if any(k in n for k in ("RECT", "RECTANGLE", "PARALLELOGRAM", "ROUNDED")):
        return "PROCESS"
    if "CYLINDER" in n:
        return "DATASTORE"
    return "SHAPE"


def _is_flow_label(text: str) -> bool:
    """Return True if the text looks like a connector flow label (Yes/No/etc.)."""
    return text.strip().lower() in {"yes", "no", "y", "n", "true", "false"}


def _is_step_number(text: str) -> bool:
    """Return True if the text looks like a step-reference number (1, 2, 3a, A, B, 5aa…)."""
    import re
    return bool(re.match(r"^[0-9A-Za-z]{1,5}$", text.strip()))


def _build_layout_summary(
    slide_number: int,
    shape_records: list[dict[str, Any]],
    *,
    connector_count: int,
    picture_count: int,
) -> str:
    if not shape_records:
        return ""

    # Separate shapes into semantic groups.
    slide_title = ""
    process_labels: list[str] = []
    decision_labels: list[str] = []
    terminal_labels: list[str] = []
    flow_labels: list[str] = []
    step_numbers: list[str] = []
    section_headers: list[str] = []
    other_labels: list[str] = []

    for record in shape_records:
        label = clean_text(record.get("text", ""))
        if not label:
            continue
        shape_type = str(record.get("shape_type", "")).upper()
        auto_name = str(record.get("auto_shape_type") or "")
        category = _auto_shape_category(auto_name)

        # Slide title placeholder
        if shape_type == "PLACEHOLDER":
            if not slide_title:
                slide_title = label
            continue

        if _is_flow_label(label):
            if label not in flow_labels:
                flow_labels.append(label)
            continue

        if _is_step_number(label):
            step_numbers.append(label)
            continue

        # Section header text boxes (short labels that precede data tables)
        if shape_type == "TEXT_BOX" and len(label) < 40:
            section_headers.append(label)
            continue

        if category == "DECISION":
            decision_labels.append(label)
        elif category in ("PROCESS", "TERMINAL", "DATASTORE"):
            process_labels.append(label) if category != "TERMINAL" else terminal_labels.append(label)
        else:
            other_labels.append(label)

    total_labeled = (
        len(process_labels) + len(decision_labels) + len(terminal_labels)
        + len(section_headers) + len(other_labels)
    )
    is_diagram = connector_count > 0 or (total_labeled >= 3 and (decision_labels or process_labels))

    if total_labeled < 2 and connector_count == 0 and picture_count == 0:
        return ""

    # Determine document type label.
    if is_diagram and (decision_labels or process_labels):
        doc_type = "WORKFLOW DIAGRAM / FLOWCHART"
    elif connector_count > 0:
        doc_type = "DIAGRAM"
    else:
        doc_type = "SLIDE"

    title_line = f'Slide {slide_number} {doc_type}'
    if slide_title:
        title_line += f': "{slide_title}"'

    parts = [title_line + "."]

    counts = []
    if process_labels:
        counts.append(f"{len(process_labels)} process step(s)")
    if decision_labels:
        counts.append(f"{len(decision_labels)} decision point(s)")
    if terminal_labels:
        counts.append(f"{len(terminal_labels)} terminal node(s)")
    if connector_count:
        counts.append(f"{connector_count} flow connector(s)")
    if picture_count:
        counts.append(f"{picture_count} embedded image(s)")
    if counts:
        parts.append("Contains: " + ", ".join(counts) + ".")

    if process_labels:
        parts.append("Process steps: " + " | ".join(process_labels) + ".")
    if terminal_labels:
        parts.append("Start/End nodes: " + " | ".join(terminal_labels) + ".")
    if decision_labels:
        parts.append("Decision points: " + " | ".join(decision_labels) + ".")
    if flow_labels:
        parts.append("Flow branch labels: " + " | ".join(flow_labels) + ".")
    if step_numbers:
        parts.append("Step reference numbers: " + " | ".join(step_numbers) + ".")
    if section_headers:
        parts.append("Section headers on slide: " + " | ".join(section_headers) + ".")
    if other_labels:
        parts.append("Other labeled elements: " + " | ".join(other_labels) + ".")

    return clean_text(" ".join(parts))


def extract_pptx(file_path: Path, rel_path: str, extract_root: Path, cfg: dict[str, Any]) -> ExtractedPPTX:
    if Presentation is None:
        raise RuntimeError("python-pptx is not installed")

    prs = Presentation(str(file_path))
    slide_width = float(getattr(prs, "slide_width", 0.0))
    slide_height = float(getattr(prs, "slide_height", 0.0))

    text_blocks: list[TextBlock] = []
    images: list[ImageItem] = []
    stats: dict[str, Any] = {
        "text_blocks": 0,
        "images_seen": 0,
        "images_kept": 0,
        "images_filtered": 0,
        "table_shapes": 0,
        "connector_shapes": 0,
        "layout_summaries": 0,
    }

    doc_order = 0
    for slide_index, slide in enumerate(prs.slides, 1):
        flat_shapes = _flatten_shapes(slide.shapes)
        flat_shapes.sort(key=lambda s: (float(getattr(s, "top", 0.0)), float(getattr(s, "left", 0.0))))
        caption_candidates: list[dict[str, Any]] = []

        for shape in flat_shapes:
            if getattr(shape, "shape_type", None) == getattr(MSO_SHAPE_TYPE, "PICTURE", None):
                continue
            if getattr(shape, "shape_type", None) == getattr(MSO_SHAPE_TYPE, "TABLE", None):
                preview_text = _table_to_text(shape, slide_index, 0)
            else:
                preview_text = _shape_text(shape)
            if not preview_text:
                continue
            caption_candidates.append(
                {
                    "page": slide_index,
                    "bbox": _shape_bbox(shape),
                    "text": preview_text,
                }
            )

        slide_blocks: list[dict[str, Any]] = []
        slide_units: list[dict[str, Any]] = []
        slide_shape_records: list[dict[str, Any]] = []
        picture_shape_counter = 0
        connector_count = 0
        table_count = 0

        for shape in flat_shapes:
            shape_type_name = _shape_type_name(shape)
            auto_shape_name = _auto_shape_name(shape)
            bbox = _shape_bbox(shape)
            rect = _shape_rect(shape)

            if getattr(shape, "shape_type", None) == getattr(MSO_SHAPE_TYPE, "PICTURE", None):
                stats["images_seen"] += 1
                picture_shape_counter += 1
                try:
                    image = shape.image
                    image_bytes = image.blob
                    ext = str(image.ext or "png").lower()
                except Exception:
                    stats["images_filtered"] += 1
                    continue

                keep = True
                reason = "ok"
                try:
                    with Image.open(BytesIO(image_bytes)) as pil_img:
                        keep, reason = is_diagram_image(pil_img.convert("RGB"), rect, slide_height, cfg)
                except Exception:
                    keep = False
                    reason = "image_open_failed"

                if not keep:
                    stats["images_filtered"] += 1
                    continue

                caption_text = find_best_caption_for_image(caption_candidates, rect, slide_width, cfg)
                ocr = compute_ocr(image_bytes)
                image_path = extract_root / "images" / rel_path / f"slide_{slide_index:03d}_img_{picture_shape_counter:03d}.{ext}"
                _save_image(image_bytes, image_path)

                slide_units.append(
                    {
                        "kind": "image",
                        "image_index": picture_shape_counter,
                        "bbox": bbox,
                        "image_path": str(image_path),
                        "ext": ext,
                        "caption_text": caption_text,
                        "filter_reason": reason,
                        "ocr_text": ocr.text,
                        "ocr_word_count": ocr.word_count,
                        "ocr_avg_conf": ocr.avg_conf,
                        "shape_name": getattr(shape, "name", None),
                    }
                )
                stats["images_kept"] += 1
                continue

            _shape_type_val = getattr(shape, "shape_type", None)
            _line_type = getattr(MSO_SHAPE_TYPE, "LINE", None)
            _connector_type = getattr(MSO_SHAPE_TYPE, "CONNECTOR", None)
            if _shape_type_val is not None and _shape_type_val in {
                t for t in [_line_type, _connector_type] if t is not None
            }:
                connector_count += 1
            if getattr(shape, "shape_type", None) == getattr(MSO_SHAPE_TYPE, "TABLE", None):
                table_count += 1
                stats["table_shapes"] += 1
                text = _table_to_text(shape, slide_index, table_count)
            else:
                text = _shape_text(shape)

            if not text:
                continue

            slide_shape_records.append(
                {
                    "text": text,
                    "shape_type": shape_type_name,
                    "auto_shape_type": auto_shape_name,
                    "bbox": bbox,
                }
            )
            slide_blocks.append(
                {
                    "id": f"S{slide_index}_T{len(slide_blocks)+1}",
                    "page": slide_index,
                    "bbox": bbox,
                    "text": text,
                    "shape_type": shape_type_name,
                    "shape_name": getattr(shape, "name", None),
                    "auto_shape_type": auto_shape_name,
                }
            )

        stats["connector_shapes"] += connector_count

        layout_summary = _build_layout_summary(
            slide_index,
            slide_shape_records,
            connector_count=connector_count,
            picture_count=picture_shape_counter,
        )
        if layout_summary:
            # y0 = -1.0 ensures the summary sorts BEFORE all slide shapes so the
            # grader sees the structural overview first, then individual content.
            slide_units.append(
                {
                    "kind": "text",
                    "id": f"S{slide_index}_SUMMARY",
                    "bbox": {"x0": 0.0, "y0": -1.0, "x1": slide_width, "y1": 0.0},
                    "text": layout_summary,
                    "shape_type": "slide_summary",
                    "shape_name": f"Slide {slide_index} Summary",
                    "auto_shape_type": None,
                }
            )
            stats["layout_summaries"] += 1

        for block in slide_blocks:
            slide_units.append({"kind": "text", **block})

        slide_units.sort(key=lambda b: (float(b["bbox"]["y0"]), float(b["bbox"]["x0"])))
        for block_index, unit in enumerate(slide_units):
            doc_order += 1
            if unit["kind"] == "text":
                text_blocks.append(
                    TextBlock(
                        block_id=str(unit["id"]),
                        page_number=slide_index,
                        block_index=block_index,
                        bbox=dict(unit["bbox"]),
                        text=str(unit["text"]),
                        sort_key=make_sort_key(slide_index, block_index),
                        document_order=doc_order,
                        shape_name=unit.get("shape_name"),
                        shape_type=unit.get("auto_shape_type") or unit.get("shape_type"),
                    )
                )
            else:
                images.append(
                    ImageItem(
                        image_index=int(unit["image_index"]),
                        page_number=slide_index,
                        block_index=block_index,
                        bbox=dict(unit["bbox"]),
                        image_path=str(unit["image_path"]),
                        ext=str(unit["ext"]),
                        caption_text=str(unit["caption_text"]),
                        filter_reason=str(unit["filter_reason"]),
                        ocr_text=str(unit["ocr_text"]),
                        ocr_word_count=int(unit["ocr_word_count"]),
                        ocr_avg_conf=float(unit["ocr_avg_conf"]),
                        sort_key=make_sort_key(slide_index, block_index),
                        document_order=doc_order,
                        shape_name=unit.get("shape_name"),
                    )
                )

    stats["text_blocks"] = len(text_blocks)
    return ExtractedPPTX(
        source_path=rel_path,
        file_type="pptx",
        slide_count=len(prs.slides),
        text_blocks=text_blocks,
        images=images,
        stats=stats,
    )


def extracted_pptx_to_jsonable(extracted: ExtractedPPTX) -> dict[str, Any]:
    return {
        "source_path": extracted.source_path,
        "file_type": extracted.file_type,
        "slide_count": extracted.slide_count,
        "stats": extracted.stats,
        "text_blocks": [asdict(x) for x in extracted.text_blocks],
        "images": [asdict(x) for x in extracted.images],
    }
