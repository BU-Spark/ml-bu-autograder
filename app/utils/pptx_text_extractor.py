from pptx import Presentation
from pathlib import Path
import json

def extract_pptx_to_chunks(pptx_path, course="CS 581", lecture_id=None):
    """
    Returns a list of RAG-ready chunks from a PowerPoint file.
    Each chunk is one slide (optionally including notes).
    """
    pptx_path = Path(pptx_path)
    prs = Presentation(pptx_path)

    if lecture_id is None:
        lecture_id = pptx_path.stem  # filename without extension

    chunks = []

    for slide_idx, slide in enumerate(prs.slides, start=1):
        # Slide title (if any)
        title_text = ""
        if slide.shapes.title and hasattr(slide.shapes.title, "text"):
            title_text = slide.shapes.title.text.strip()

        # Body text (all text shapes except title)
        body_parts = []
        for shape in slide.shapes:
            if shape == slide.shapes.title:
                continue
            if not hasattr(shape, "has_text_frame") or not shape.has_text_frame:
                continue
            text = shape.text_frame.text.strip()
            if text:
                body_parts.append(text)

        body_text = "\n".join(body_parts).strip()

        # Speaker notes (if present)
        notes_text = ""
        if getattr(slide, "has_notes_slide", False) and slide.notes_slide:
            ntf = getattr(slide.notes_slide, "notes_text_frame", None)
            if ntf:
                notes_lines = [p.text for p in ntf.paragraphs if p.text.strip()]
                notes_text = "\n".join(notes_lines).strip()

        # Skip completely empty slides
        if not (title_text or body_text or notes_text):
            continue

        # Combine into a single text field for RAG
        pieces = []
        if title_text:
            pieces.append(f"Slide title: {title_text}")
        if body_text:
            pieces.append(body_text)
        if notes_text:
            pieces.append(f"Speaker notes:\n{notes_text}")

        combined_text = "\n\n".join(pieces).strip()

        chunk = {
            "id": f"{lecture_id}_slide_{slide_idx}",
            "text": combined_text,
            "metadata": {
                "course": course,
                "lecture_id": lecture_id,
                "source_file": pptx_path.name,
                "slide_index": slide_idx,
                "modality": "pptx_slide",
            },
        }
        chunks.append(chunk)

    return chunks